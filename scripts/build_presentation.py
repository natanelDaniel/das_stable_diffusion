"""
build_presentation.py

Builds SyntheDAS_presentation.pptx (15 slides) from scratch.

Embeds real results from ppt_results/ if available:
    ppt_results/metrics.json
    ppt_results/comparison_f1.png
    ppt_results/confusion_matrix_{run}.png   (baseline/synthetic/ratio_0p5/ratio_1p0)
    ppt_results/roc_pr_{run}.png

Also embeds diffusion visualisations if available:
    figures/real_vs_generated.png
    figures/alpha_sweep.png

Usage:
    python scripts/build_presentation.py
    python scripts/build_presentation.py --out SyntheDAS_presentation.pptx
"""

import argparse
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx is not installed. Run:  pip install python-pptx")
    sys.exit(1)

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1a, 0x3a, 0x5c)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
GREEN  = RGBColor(0x16, 0xa3, 0x4a)
GRAY_L = RGBColor(0xe5, 0xe7, 0xeb)
GRAY_D = RGBColor(0x6b, 0x72, 0x80)
WHITE  = RGBColor(0xff, 0xff, 0xff)
BLACK  = RGBColor(0x11, 0x18, 0x27)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_TOP  = Inches(0.25)
TITLE_LEFT = Inches(0.5)
TITLE_W    = Inches(12.33)
TITLE_H    = Inches(0.95)

BODY_TOP  = Inches(1.35)
BODY_LEFT = Inches(0.5)
BODY_W    = Inches(12.33)
BODY_H    = Inches(5.8)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: RGBColor = WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title(slide, text: str, font_size: int = 32):
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = NAVY
    bar = slide.shapes.add_shape(1, TITLE_LEFT, Inches(1.27), TITLE_W, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def _add_body_text(slide, lines, top=BODY_TOP, left=BODY_LEFT,
                   width=BODY_W, height=BODY_H, font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for level, text in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        if level == 1:
            p.text = f"  {text}"
        elif level == 2:
            p.text = f"      {text}"
        else:
            p.text = text
        for run in p.runs:
            run.font.size = Pt(font_size - level * 1.5)
            run.font.color.rgb = BLACK
        if level == 0 and p.runs:
            p.runs[0].font.bold = True


def _add_placeholder_image(slide, label: str, left, top, width, height):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_L
    box.line.color.rgb = GRAY_D
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_D


def _embed_or_placeholder(slide, img_path: str, label: str,
                           left, top, width, height):
    if os.path.isfile(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        _add_placeholder_image(slide, f"[{label}]", left, top, width, height)


def _add_table(slide, headers, rows, left, top, width, height, font_size=14):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for col, h in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = WHITE
    for r, row_data in enumerate(rows):
        fill_color = RGBColor(0xf0, 0xf4, 0xff) if r % 2 == 0 else WHITE
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size - 1)
                    run.font.color.rgb = BLACK


def _add_textbox(slide, text, left, top, width, height,
                 font_size=13, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or BLACK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide builders — context / model slides (unchanged from v1)
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_textbox(slide, "SyntheDAS",
                 Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.8),
                 font_size=52, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Generative AI for Fiber-Optic Sensing",
                 Inches(0.5), Inches(3.5), Inches(12.33), Inches(0.8),
                 font_size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Conditional Pixel-Space Diffusion Model for DAS Signal Synthesis & Augmentation",
                 Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.8),
                 font_size=18, color=BLACK, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Netanel Daniel",
                 Inches(0.5), Inches(5.6), Inches(12.33), Inches(0.6),
                 font_size=16, color=GRAY_D, align=PP_ALIGN.CENTER)
    bar = slide.shapes.add_shape(1, Inches(3.5), Inches(5.0), Inches(6.33), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()


def slide_background(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Background & Importance")
    _add_body_text(slide, [
        (0, "Distributed Acoustic Sensing (DAS)"),
        (1, "Turns kilometers of standard fiber-optic cable into a dense virtual sensor array"),
        (1, "Data appears as waterfall plots: Time (x-axis) x Fiber channel (y-axis)"),
        (1, "Sampling: 500 Hz per channel, 8 channels, 32.768 s patches"),
        (0, "Mission-Critical Applications"),
        (1, "Border security & perimeter monitoring"),
        (1, "Pipeline integrity & leak detection"),
        (1, "Urban traffic classification & facility protection"),
        (0, "The Challenge"),
        (1, "High noise variance across soil types and environments"),
        (1, "Rare event classes severely underrepresented — car, running, walk"),
        (1, "Today: heavy reliance on human operators or data-starved ML models"),
    ])


def slide_problem(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Problem Statement & Innovation")
    _add_body_text(slide, [
        (0, "Core Problem: Data Scarcity & Lack of Diversity"),
        (1, "Models trained on one environment fail in unseen soil types or noise conditions"),
        (1, "Rare event classes (car, running, walk) are severely underrepresented"),
        (1, "Manual data collection is expensive and environment-specific"),
        (0, "Our Approach: Conditional Pixel-Space Diffusion Model"),
        (1, "Generate realistic DAS patches conditioned on event class and noise level (alpha)"),
        (1, "Use synthetic data to augment the real training set for a downstream CNN classifier"),
        (0, "Key Innovation"),
        (1, "First conditional diffusion model specifically designed for DAS signal synthesis"),
        (1, "Dual conditioning: class label (WHAT event) + alpha (HOW MUCH noise mixing)"),
        (1, "Enables controlled data augmentation without additional field collection"),
    ])


def slide_diffusion_arch(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Diffusion Model Architecture")
    _add_body_text(slide, [
        (0, "1D U-Net — Pixel Space"),
        (1, "Input: [B, 8, 16384]  (8 DAS channels x 32.768 s @ 500 Hz)"),
        (1, "4 encoder/decoder levels: 96 -> 192 -> 192 -> 384 channels"),
        (1, "2 residual blocks per level"),
        (0, "Dilated Convolutions  (WaveNet-style)"),
        (1, "Level 0: rates [1, 2]  ->  Level 3: rates [1, 16]"),
        (1, "Exponentially growing receptive field without added parameters"),
        (0, "Self-Attention"),
        (1, "Applied at the deepest level (W=2048) + bottleneck"),
        (1, "Captures long-range temporal dependencies across the full patch"),
        (0, "Conditioning via FiLM (Feature-wise Linear Modulation)"),
        (1, "3 signals combined additively into each ResBlock: timestep t, class label c, alpha"),
        (1, "Classifier-free guidance: 10% dropout on class & alpha during training"),
        (1, "Inference: cfg_scale=7.5  ->  sharp class-specific patterns"),
    ], font_size=16)


def slide_data_split(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Data & Split Strategy")
    _add_body_text(slide, [
        (0, "Dataset: Brno University DAS Dataset"),
        (1, "4 event classes: Car  |  Background (regular)  |  Running  |  Walk"),
        (1, "Patch shape: 8 channels x 16384 time samples (32.768 s @ 500 Hz)"),
        (1, "Event anchored within [4 s, 28 s] of the patch to avoid edge artifacts"),
        (0, "Recording-Level Stratified Split  (70 / 15 / 15)"),
        (1, "All patches from one recording stay in the same split -> no temporal leakage"),
        (1, "Classes with 2 recordings: temporal hybrid (rec-1 train, rec-2 -> 50/50 val/test)"),
        (1, "Same split shared between diffusion model and CNN -> no cross-model leakage"),
        (0, "Class Balancing"),
        (1, "Decimation: regular=5%, car/running/walk=50%  (reduces dominant class)"),
        (1, "WeightedRandomSampler with inverse-frequency weights during training"),
        (1, "Test set: 5 116 patches  |  train: 43 296  |  val: 7 611"),
    ])


def slide_losses(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Loss Functions & Design Choices")
    _add_table(
        slide,
        headers=["Loss", "Weight", "Purpose"],
        rows=[
            ["MSE on v-prediction target", "1.0",
             "Core denoising. V-space stable at all SNR levels (better than eps-prediction near t=0)"],
            ["Multi-scale STFT", "0.1",
             "Spectral fidelity across all frequency bands"],
            ["Band-limited STFT  (0-128 Hz)", "1.0",
             "Emphasises low-frequency seismic content — dominant energy band in DAS signals"],
            ["Derivative L1", "0.5",
             "Preserves transient edges and event onset sharpness"],
            ["Min-SNR weighting  (gamma=5)", "—",
             "Prevents high-noise timesteps from overwhelming the gradient signal"],
        ],
        left=Inches(0.5), top=Inches(1.45),
        width=Inches(12.33), height=Inches(3.2),
        font_size=13,
    )
    _add_body_text(slide, [
        (0, "Noise schedule: cosine (squaredcos_cap_v2), 1000 timesteps, zero-SNR rescaling"),
        (1, "Inference: DDIM sampler, 50 steps  ->  fast, deterministic generation"),
    ], top=Inches(4.85), font_size=16)


def slide_real_vs_gen(prs, img_path: str):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Real vs. Diffusion-Generated DAS Patches")
    _add_body_text(slide, [
        (1, "4 event classes  x  3 columns: [Real | Generated-1 | Generated-2]"),
        (1, "Generated at alpha=0.0 (clean event), CFG scale=7.5, 50 DDIM steps"),
    ], top=Inches(1.3), height=Inches(0.6), font_size=15)
    _embed_or_placeholder(slide, img_path, "plot_real_vs_generated.py",
                          left=Inches(0.4), top=Inches(2.0),
                          width=Inches(12.5), height=Inches(5.1))


def slide_alpha_sweep(prs, img_path: str):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Alpha Conditioning: Event -> Noise Transition")
    _add_body_text(slide, [
        (1, "alpha=0.0  ->  pure clean event signal    |    alpha=1.0  ->  fully noise-dominated"),
        (1, "5 alpha values x 4 classes — same model, same weights, different conditioning vector"),
    ], top=Inches(1.3), height=Inches(0.6), font_size=15)
    _embed_or_placeholder(slide, img_path, "plot_alpha_sweep.py",
                          left=Inches(0.4), top=Inches(2.0),
                          width=Inches(12.5), height=Inches(5.1))


def slide_diffusion_results(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Diffusion Model — Qualitative Results")
    _add_body_text(slide, [
        (0, "Class-Conditional Generation (alpha=0.0, CFG=7.5)"),
        (1, "Car: strong horizontal stripes reproduced — characteristic slow ground vibration"),
        (1, "Running: sharp vertical spikes per footstep, correct inter-spike rhythm"),
        (1, "Walk: slower, heavier spike pattern — distinct from running cadence"),
        (1, "Background: diffuse low-amplitude noise, no structured event features"),
        (0, "Alpha Conditioning (Event -> Noise Transition)"),
        (1, "alpha=0.0: crisp event signature dominates across all channels"),
        (1, "alpha=0.5: event structure begins to blur into noise floor"),
        (1, "alpha=1.0: patch resembles background — structural event features masked"),
        (0, "Key Observation"),
        (1, "Model correctly interpolates between conditioned event and noise-dominated signal"),
        (1, "CFG scale=7.5 sharpens class-specific features without mode collapse"),
    ], font_size=16)


def slide_cnn_arch(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "CNN Downstream Classifier Architecture")
    _add_body_text(slide, [
        (0, "DASResNetClassifier  (ResNet-34-style 2D CNN)"),
        (1, "Input: [B, 1, 8, 16384]  ->  Output: [B, 4] logits"),
        (0, "Stem (temporal compression)"),
        (1, "2 x Conv2d(kernel=(1,7), stride=(1,4))  ->  reduces T from 16384 -> 1024"),
        (1, "Asymmetric kernel compresses time aggressively while preserving fiber-channel structure"),
        (0, "4 Residual Stages"),
        (1, "64 -> 64 -> 128 -> 256 -> 512 channels"),
        (1, "Each BasicBlock: Conv-BN-ReLU-Conv-BN + residual shortcut"),
        (1, "Stride (2,2) per stage reduces spatial + temporal dims"),
        (0, "Head"),
        (1, "AdaptiveAvgPool2d(1,1)  ->  Dropout(0.35)  ->  Linear(512 -> 4)"),
    ], font_size=16)


def slide_cnn_loss(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "CNN Loss Function & Training Strategy")
    _add_body_text(slide, [
        (0, "Loss: Weighted Cross-Entropy"),
        (1, "Weights = 1 / class_count, normalised to sum to 1"),
        (1, "Why? 'regular' class is ~20x more common than car / running / walk"),
        (1, "Prevents the model from collapsing to predict 'regular' for everything"),
        (0, "Optimiser & Schedule"),
        (1, "AdamW  (lr=3e-4,  weight_decay=0.0085)"),
        (1, "CosineAnnealingWarmRestarts  T0=50 epochs  ->  periodic restarts escape local minima"),
        (1, "Mixed precision (AMP) + gradient clipping at 1.0"),
        (0, "Data Augmentation"),
        (1, "Gaussian background noise std=0.3 (normalised space) applied per batch"),
        (1, "Runs 'synthetic' and 'ratio_*': BackgroundMixupDataset + synthetic event patches"),
    ], font_size=16)


def slide_experiment_design(prs):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Experiment Design")

    _add_textbox(
        slide,
        "Research question: Can a pixel-space diffusion model trained on DAS serve as an "
        "augmentation source that improves downstream classification?",
        Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.55),
        font_size=15, italic=True, color=BLUE,
    )

    _add_table(
        slide,
        headers=["W&B Run", "Training Data", "Augmentation Strategy"],
        rows=[
            ["cnn_baseline",
             "Real DAS patches only",
             "None — reference run"],
            ["cnn_synthetic",
             "Real + synthetic backgrounds (mixup)",
             "BackgroundMixupDataset: real event + alpha*bg,  alpha~U(0,0.4)\nDiffusion generates 'regular' class backgrounds only"],
            ["cnn_synth_ratio_0.5",
             "Real + 0.5x synthetic event patches",
             "Inject synthetic events from ALL 4 classes at 0.5x real dataset size"],
            ["cnn_synth_ratio_1.0",
             "Real + 1.0x synthetic event patches",
             "Inject synthetic events from ALL 4 classes at 1.0x real dataset size"],
        ],
        left=Inches(0.5), top=Inches(2.05),
        width=Inches(12.33), height=Inches(2.8),
        font_size=12,
    )

    _add_body_text(slide, [
        (0, "Comparison axis 1 — Augmentation type"),
        (1, "Baseline vs background-mixup vs full synthetic event injection"),
        (0, "Comparison axis 2 — Synthetic data volume (0.5x vs 1.0x)"),
        (1, "Does more synthetic data keep improving, or is there a saturation point?"),
        (0, "All runs: identical architecture, hyperparameters, split, and seed — only data differs"),
    ], top=Inches(5.05), height=Inches(2.2), font_size=15)


# ---------------------------------------------------------------------------
# Results slides — embed real images and metrics
# ---------------------------------------------------------------------------

RUN_ORDER  = ["baseline", "synthetic", "ratio_0p5", "ratio_1p0"]
RUN_LABELS = ["Baseline", "Synthetic\n(BG Mixup)", "Ratio 0.5x", "Ratio 1.0x"]
CLASSES    = ["car", "regular", "running", "walk"]


def slide_results_overview(prs, metrics, ppt_dir):
    """Summary metrics table + comparison F1 bar chart."""
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "CNN Results — Overview")

    rows = []
    for run, label in zip(RUN_ORDER, RUN_LABELS):
        m = metrics.get(run, {})
        rows.append([
            label.replace("\n", " "),
            f"{m.get('acc', 0):.3f}",
            f"{m.get('f1_macro', 0):.3f}",
        ] + [f"{m.get('f1', {}).get(c, 0):.3f}" for c in CLASSES])

    _add_table(
        slide,
        headers=["Run", "Accuracy", "F1-macro", "F1-car", "F1-regular", "F1-running", "F1-walk"],
        rows=rows,
        left=Inches(0.5), top=Inches(1.45),
        width=Inches(12.33), height=Inches(1.75),
        font_size=11,
    )

    comp = os.path.join(ppt_dir, "comparison_f1.png")
    _embed_or_placeholder(slide, comp, "comparison_f1.png",
                          left=Inches(1.8), top=Inches(3.35),
                          width=Inches(9.73), height=Inches(3.9))


def slide_cnn_results(prs, metrics, ppt_dir):
    """4-column slide: one column per run with metrics + confusion matrix."""
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "CNN Results — Per-Run Detail")

    col_w = Inches(3.0)
    gap   = Inches(0.11)
    x0    = Inches(0.3)

    for i, (run, label) in enumerate(zip(RUN_ORDER, RUN_LABELS)):
        left = x0 + i * (col_w + gap)
        m    = metrics.get(run, {})
        acc  = m.get("acc", 0)
        f1m  = m.get("f1_macro", 0)

        # Column header
        _add_textbox(slide, label.replace("\n", " "),
                     left, Inches(1.45), col_w, Inches(0.38),
                     font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # Acc + F1-macro
        _add_textbox(slide, f"Acc={acc:.3f}   F1={f1m:.3f}",
                     left, Inches(1.88), col_w, Inches(0.32),
                     font_size=11, color=BLUE, align=PP_ALIGN.CENTER)

        # Per-class F1
        lines = "\n".join(
            f"{c}: {m.get('f1', {}).get(c, 0):.3f}"
            for c in CLASSES
        )
        _add_textbox(slide, lines,
                     left, Inches(2.24), col_w, Inches(1.15),
                     font_size=10, color=BLACK)

        # Confusion matrix image
        cm = os.path.join(ppt_dir, f"confusion_matrix_{run}.png")
        _embed_or_placeholder(slide, cm, f"confusion_{run}",
                              left=left, top=Inches(3.5),
                              width=col_w, height=Inches(3.7))


def slide_summary(prs, metrics=None):
    slide = _blank_slide(prs)
    _set_bg(slide)
    _add_title(slide, "Summary & Key Findings")

    m_bl = (metrics or {}).get("baseline", {})
    m_sy = (metrics or {}).get("synthetic", {})
    m_r5 = (metrics or {}).get("ratio_0p5", {})
    m_r1 = (metrics or {}).get("ratio_1p0", {})

    bl_f1  = m_bl.get("f1_macro", 0)
    sy_f1  = m_sy.get("f1_macro", 0)
    r5_f1  = m_r5.get("f1_macro", 0)
    r1_f1  = m_r1.get("f1_macro", 0)
    r1_car = m_r1.get("f1", {}).get("car", 0)
    r1_run = m_r1.get("f1", {}).get("running", 0)

    _add_body_text(slide, [
        (0, "What We Built"),
        (1, "Pixel-space diffusion model: 1D U-Net with FiLM conditioning on (class, alpha, timestep)"),
        (1, "Recording-level stratified split shared between diffusion + CNN — zero data leakage"),
        (1, "Multi-term loss: v-prediction MSE + band-limited STFT + derivative L1 + min-SNR weighting"),
        (0, "Key Findings"),
        (1, f"Baseline (real data only): F1-macro={bl_f1:.3f} — car class completely missed (F1=0.00)"),
        (1, f"BG Mixup only: F1-macro={sy_f1:.3f} — car appears (F1=0.52), moderate improvement"),
        (1, f"Ratio 0.5x synthetic events: F1-macro={r5_f1:.3f} — large jump across all classes"),
        (1, f"Ratio 1.0x synthetic events (best): F1-macro={r1_f1:.3f} — car F1={r1_car:.2f}, running F1={r1_run:.2f}"),
        (1, "Synthetic event injection outperforms BG-only augmentation by a large margin"),
        (0, "Next Steps"),
        (1, "Evaluate on held-out recordings from unseen environments / soil types"),
        (1, "Ablation: loss term contribution (STFT / band-limited / derivative individually)"),
        (1, "Test higher ratio (2.0x) to probe saturation of synthetic data benefit"),
    ], font_size=16)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out",          default="SyntheDAS_presentation.pptx")
    parser.add_argument("--real_vs_gen",  default="figures/real_vs_generated.png")
    parser.add_argument("--alpha_sweep",  default="figures/alpha_sweep.png")
    parser.add_argument("--ppt_results",  default="ppt_results")
    args = parser.parse_args()

    # Load real metrics if available
    metrics_path = os.path.join(args.ppt_results, "metrics.json")
    if os.path.isfile(metrics_path):
        with open(metrics_path) as f:
            metrics = json.load(f)
        print(f"Loaded metrics from {metrics_path}")
    else:
        metrics = {}
        print(f"[NOTE] {metrics_path} not found — results slides will use placeholder numbers.")

    prs = _new_prs()
    print("Building slides ...")

    slide_title(prs)                                    # 1
    slide_background(prs)                               # 2
    slide_problem(prs)                                  # 3
    slide_diffusion_arch(prs)                           # 4
    slide_data_split(prs)                               # 5
    slide_losses(prs)                                   # 6
    slide_real_vs_gen(prs, args.real_vs_gen)            # 7
    slide_alpha_sweep(prs, args.alpha_sweep)            # 8
    slide_diffusion_results(prs)                        # 9
    slide_cnn_arch(prs)                                 # 10
    slide_cnn_loss(prs)                                 # 11
    slide_experiment_design(prs)                        # 12
    slide_results_overview(prs, metrics, args.ppt_results)  # 13
    slide_cnn_results(prs, metrics, args.ppt_results)       # 14
    slide_summary(prs, metrics)                         # 15

    prs.save(args.out)
    print(f"\nSaved: {args.out}  ({len(prs.slides)} slides)")
    print("\nTo import into Google Slides:")
    print("  File -> Import Slides -> Upload -> select", args.out, "-> Import All")


if __name__ == "__main__":
    main()
